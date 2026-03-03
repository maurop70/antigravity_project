import os
import glob
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Directory containing this script (and presumably the assets if moved, but we will target the artifact dir)
# Adjust these paths as necessary based on where the images are actually stored.
# Based on the user session, they are in the brain artifact dir.
ARTIFACT_DIR = r"C:\Users\mpetr\.gemini\antigravity\brain\b8bd6a3b-8b45-4768-92e2-aff15277836d"

def find_image(prefix):
    """Finds the first image file starting with specific prefix in the artifact dir."""
    pattern = os.path.join(ARTIFACT_DIR, f"{prefix}*.png")
    files = glob.glob(pattern)
    if files:
        # Sort by modification time to get the latest if multiple exist
        files.sort(key=os.path.getmtime, reverse=True)
        return files[0]
    return None

def create_presentation():
    prs = Presentation()

    # Define Slide Layouts
    # 0 = Title, 1 = Title + Content, 5 = Title Only (good for custom images), 6 = Blank
    slide_layout_title = prs.slide_layouts[0]
    slide_layout_content = prs.slide_layouts[1]
    slide_layout_blank = prs.slide_layouts[6]

    # --- Slide 1: Title ---
    slide1 = prs.slides.add_slide(slide_layout_title)
    title1 = slide1.shapes.title
    subtitle1 = slide1.placeholders[1]
    title1.text = "Vegan Coffee Gelato: A Modern Indulgence"
    subtitle1.text = "From Artisan Craft to Industrial Scale"
    
    # Add Hero Image
    img_path = find_image("slide1_hero")
    if img_path:
        # Add image to background or side. Let's make it a nice side element or background.
        # For simplicity in this script, we'll place it at the bottom.
        left = Inches(1)
        top = Inches(3.5)
        height = Inches(3.5)
        slide1.shapes.add_picture(img_path, left, top, height=height)

    # --- Slide 2: Small-Scale Preparation ---
    slide2 = prs.slides.add_slide(slide_layout_content)
    title2 = slide2.shapes.title
    title2.text = "Small-Scale Preparation"
    content2 = slide2.placeholders[1]
    content2.text = (
        "• Artisan Batch Method\n"
        "• Cold brew infusion for 18 hours\n"
        "• Macadamia milk base\n"
        "• Churned at -8°C for optimal structure"
    )
    
    img_path = find_image("slide2_artisan")
    if img_path:
        left = Inches(5.5)
        top = Inches(2)
        height = Inches(4)
        slide2.shapes.add_picture(img_path, left, top, height=height)

    # --- Slide 3: Industrial Manufacturing ---
    slide3 = prs.slides.add_slide(slide_layout_content)
    title3 = slide3.shapes.title
    title3.text = "Industrial Manufacturing Process"
    content3 = slide3.placeholders[1]
    content3.text = (
        "• Scalable production without compromising quality\n"
        "• Continuous freezers with precision overrun control (30-40%)\n"
        "• Shock freezing tunnel (-40°C) to lock in texture\n"
        "• Automated filling and capping"
    )
    
    img_path = find_image("slide3_industrial")
    if img_path:
        left = Inches(5.5)
        top = Inches(2)
        height = Inches(4)
        slide3.shapes.add_picture(img_path, left, top, height=height)

    # --- Slide 4: Flavor & Aromatic Profile ---
    slide4 = prs.slides.add_slide(slide_layout_content)
    title4 = slide4.shapes.title
    title4.text = "Flavor and Aromatic Profile"
    content4 = slide4.placeholders[1]
    content4.text = (
        "• Primary: Deep roasted Arabica coffee notes\n"
        "• Secondary: Sweet caramel and nut undertones\n"
        "• Mouthfeel: Rich and creamy, mimicking dairy fat\n"
        "• Clean finish with no bitterness"
    )
    
    img_path = find_image("slide4_flavor")
    if img_path:
        left = Inches(5.5)
        top = Inches(2)
        height = Inches(4)
        slide4.shapes.add_picture(img_path, left, top, height=height)

    # --- Slide 5: Texture & Sensory Experience ---
    slide5 = prs.slides.add_slide(slide_layout_content)
    title5 = slide5.shapes.title
    title5.text = "Texture and Sensory Experience"
    content5 = slide5.placeholders[1]
    content5.text = (
        "• Smoothness: < 20 micron ice crystal size\n"
        "• Scoopability: Perfect at -14°C serving temp\n"
        "• Melting: Slow, consistent melt down\n"
        "• Body: Chewy yet light (high solids, moderate air)"
    )
    
    img_path = find_image("slide5_texture")
    if img_path:
        left = Inches(5.5)
        top = Inches(2)
        height = Inches(4)
        slide5.shapes.add_picture(img_path, left, top, height=height)

    # --- Slide 6: Technical Specifications ---
    slide6 = prs.slides.add_slide(slide_layout_content)
    title6 = slide6.shapes.title
    title6.text = "Technical Specifications"
    content6 = slide6.placeholders[1]
    # Creating a table would be ideal, but text list is safer for layout
    content6.text = (
        "Parameter            Value\n"
        "-----------------------------------\n"
        "Total Solids         38 - 40%\n"
        "Fat                  6 - 8% (Plant-based)\n"
        "Protein              2 - 3%\n"
        "Sugars               22 - 24%\n"
        "PAC (Anti-Freezing)  270\n"
        "POD (Sweetness)      180"
    )

    # Save
    output_path = os.path.join(ARTIFACT_DIR, "Vegan_Coffee_Gelato.pptx")
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_presentation()
